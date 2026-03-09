"""
Build Alabama Tree Cover GEE presentation as PowerPoint (.pptx).
Requires: pip install python-pptx
Run: python build_presentation.py
Output: Alabama_Tree_Cover_GEE_Presentation.pptx in the project root.
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    raise

# ----- Slide content (title, body bullets, optional speaker note) -----
SLIDES = [
    {
        "title": "Alabama Tree Cover Mapping on Google Earth Engine",
        "subtitle": "State-scale wooded-area map with free data and a scalable, reproducible pipeline",
        "body": [],
        "note": "One-sentence pitch: We built a statewide forest/non-forest map for Alabama on GEE using free optical + radar data and NLCD-derived labels, with no manual labeling and a path to temporal updates (e.g., 2023).",
    },
    {
        "title": "Problem & Goal",
        "body": [
            "Problem:",
            "• Need state-scale tree/wooded cover for Alabama; free data, GEE, scalable (no thousands of manual downloads)",
            "• No manual polygon labeling at scale",
            "",
            "Goal:",
            "• State-scale wooded-area map (forest vs non-forest) at 30 m or 10 m, for a chosen year (e.g., 2023)",
            "• Reproducible pipeline: ingest → composite/features → train → predict → export",
            "• Temporally updated relative to NLCD alone (NLCD 2020/2021 vs our composite year)",
        ],
        "note": "Emphasize one root, many paths—we evaluated options and chose the best trade-off.",
    },
    {
        "title": "Design Approach: Fractal Possibilities, Optimal Path",
        "body": [
            "We did not pick a single recipe by default.",
            "",
            "• Root: Alabama state-scale tree cover on GEE; free data; no manual labeling; lightweight",
            "• Fractal: Many valid paths (optical only, optical + radar, temporal, GEDI height, index-threshold, etc.)",
            "• Optimal path: We evaluated candidate paths against clear criteria and selected the approach that best balances:",
            "  Optimization, scalability, pipeline efficiency, speed, accuracy, data availability",
        ],
        "note": "This framing shows the work is methodical and criteria-driven, not ad hoc.",
    },
    {
        "title": "Optimization Criteria (How We Chose the Path)",
        "body": [
            "Optimization  – Tunable pipeline (composite, features, model); not ad hoc",
            "Scalability   – Full Alabama; tiled export; batch; no per-scene manual steps",
            "Data pipeline – Clear flow: ingest → composite/features → train → predict → export; GEE-native",
            "Efficiency    – Bounded compute and I/O",
            "Speed         – Time to first result and full-state result tractable (hours to days)",
            "Accuracy      – Measurable vs reference (NLCD, FIA); labels support validation",
            "Data availability – Free, Alabama coverage, ideally CONUS/global",
        ],
        "note": "Path A (Optical + NLCD → RF) scored best on pipeline, scalability, accuracy, and data availability.",
    },
    {
        "title": "Candidate Paths We Evaluated",
        "body": [
            "A (chosen): Optical (S2) + terrain → NLCD forest classes → RF in GEE; 30 m export",
            "B: Optical + terrain → GEDI height → RF; GEDI sparse → scalability issues",
            "C: Temporal optical + terrain → NLCD → RF; more signal, more compute",
            "D: Single composite + terrain → NLCD/FIA validation only → index-threshold; fast but weak accuracy",
            "",
            "Why Path A: Best balance of pipeline clarity, scalability, accuracy, data availability; no GEDI grid aggregation.",
        ],
        "note": "Path C is the natural upgrade for more accuracy; Path B for validation or height product.",
    },
    {
        "title": "What We Built — Implemented Pipeline (Path A+)",
        "body": [
            "Product: Wooded-area map (0/1) for Alabama at 10 m; target year set by composite (e.g., 2023)",
            "Labels: NLCD 2021 — forest = 41, 42, 43; optional NAIP-derived points for better accuracy",
            "Predictors:",
            "  • Sentinel-2 seasonal composites (leaf-on, leaf-off): bands + NDVI, NDMI, NBR, red-edge NDVI",
            "  • Sentinel-1 seasonal composites (VV, VH in dB; ratios)",
            "  • Slope (SRTM)",
            "Method: Random Forest in GEE → classify full state → export GeoTIFF (probability + binary @ 0.5)",
        ],
        "note": "Implementation includes radar and seasonal composites for better forest vs crops/grass separation.",
    },
    {
        "title": "Pipeline Architecture",
        "body": [
            "Sentinel-2 → Cloud mask → Seasonal composites (leaf-on/leaf-off) → Indices (NDVI, NDMI, NBR, RENDVI)",
            "Sentinel-1 → Seasonal composites (VV, VH dB) → Ratios",
            "SRTM → Slope",
            "NLCD 2021 → Forest mask (41,42,43 → 1) → Training labels",
            "",
            "Stratified sampling → Train RF (GEE) → Classify full state → Export to Drive (GeoTIFF)",
            "Outputs: probability map + binary forest @ 0.5 + NLCD baseline for validation",
        ],
        "note": "Everything runs in GEE; export is tiled/batch via Tasks to avoid memory limits.",
    },
    {
        "title": "Why This Is Useful vs. Using NLCD Alone",
        "body": [
            "Temporal: NLCD = 2020/2021 release cycle → Our map = composite year (e.g., 2023), temporally updated",
            "Signal: USGS workflow → Our features (S2 + S1 seasonal + indices + slope); can differ where land cover changed",
            "Resolution: 30 m → 30 m or 10 m export; same comparison baseline",
            "Labels: NLCD used as training source; model transfers to new imagery",
        ],
        "note": "We use NLCD as label source and produce an updated map for a chosen year.",
    },
    {
        "title": "Technical Implementation Summary",
        "body": [
            "Platform: Google Earth Engine (Code Editor JS + optional Python/geemap)",
            "Training: Stratified sampling; ~2k–6k points per class; 300 trees; in-GEE RF",
            "Memory: No full-state layers on map; export via Tasks; evaluation on stratified sample or uploaded assets",
            "Exports: (1) Probability of forest, (2) Binary forest @ 0.5, (3) NLCD forest mask (30 m) for QGIS validation",
            "Validation: In-GEE confusion matrix; gee_rf_accuracy_from_assets.js; QGIS workflows in METRICS_IN_QGIS.md",
        ],
        "note": "Key artifacts: gee_alabama_tree_cover_rf.js, gee_rf_accuracy_from_assets.js, gee_explore_alabama_geemap.py, METRICS_IN_QGIS.md.",
    },
    {
        "title": "Validation & Metrics",
        "body": [
            "In-GEE: Confusion matrix (RF vs NLCD) on stratified sample; accuracy, kappa, producer/consumer accuracy in Console",
            "From assets: Upload RF + NLCD GeoTIFFs → gee_rf_accuracy_from_assets.js → multi-seed sampling, optional CSV",
            "In QGIS: SAGA Confusion Matrix or Raster Calculator agreement map; zonal statistics by county/watershed for area and % forest",
        ],
        "note": "We can report RF vs NLCD accuracy (stratified) and metrics in QGIS; FIA can be added for independent validation.",
    },
    {
        "title": "Optional Upgrade — NAIP-Derived Labels",
        "body": [
            "Limitation: Training on NLCD approximates NLCD; edge and class errors carry over",
            "Upgrade: Manual 0/1 labels from NAIP imagery",
            "  • gee_create_naip_label_points.js exports point grid (core_forest, core_nonforest, edge)",
            "  • Label in QGIS with NAIP basemap; fill forest_label (0/1)",
            "  • Upload as asset; set LABEL_MODE = 'NAIP_POINTS' in main script",
            "Result: Same pipeline, better label quality, no change to export or scale",
        ],
        "note": "Optional follow-up for higher accuracy without changing the rest of the pipeline.",
    },
    {
        "title": "Deliverables & Repo Structure",
        "body": [
            "Product: Alabama wooded-area map (0/1) at 10 m + NLCD baseline at 30 m for chosen year",
            "Code: gee_alabama_tree_cover_rf.js (main pipeline), gee_rf_accuracy_from_assets.js, gee_create_naip_label_points.js, gee_explore_alabama_geemap.py",
            "Docs: PROJECT_SPEC.md, PROJECT_BOUNDARY.md, PATH_ANALYSIS.md, README.md, METRICS_IN_QGIS.md",
        ],
        "note": "Everything is documented and runnable; no proprietary or paywalled data.",
    },
    {
        "title": "Summary",
        "body": [
            "Goal: State-scale tree cover for Alabama on GEE with free data and a scalable pipeline",
            "Approach: Evaluated paths against criteria; selected Path A and extended with Sentinel-1 and seasonal composites",
            "Output: Wooded-area map (forest/non-forest) at 10 m for a chosen year; temporally updated vs NLCD; validated in GEE and QGIS",
            "Next steps (optional): Path C (temporal) for more signal; NAIP labels for better accuracy; GEDI for validation or height product",
        ],
        "note": "Close with: reproducible, criteria-driven, scalable Alabama tree cover on GEE with no manual labeling at scale.",
    },
    {
        "title": "Thank you",
        "subtitle": "Questions?",
        "body": [
            "Repo / code: See README.md and PROJECT_SPEC.md",
            "Validation: METRICS_IN_QGIS.md, gee_rf_accuracy_from_assets.js",
        ],
        "note": "",
    },
]


def add_slide(prs, slide_spec):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    # Title
    left, top, width = Inches(0.5), Inches(0.4), Inches(9)
    tx = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = slide_spec["title"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x1a, 0x47, 0x2e)

    # Subtitle (if present)
    if slide_spec.get("subtitle"):
        tx2 = slide.shapes.add_textbox(left, Inches(1.0), width, Inches(0.5))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = slide_spec["subtitle"]
        p2.font.size = Pt(16)
        p2.font.italic = True
        p2.font.color.rgb = RgbColor(0x4a, 0x4a, 0x4a)
        body_top = 1.6
    else:
        body_top = 1.2

    # Body bullets
    body = slide_spec.get("body", [])
    if body:
        tx_body = slide.shapes.add_textbox(left, Inches(body_top), width, Inches(5.5))
        tf_body = tx_body.text_frame
        tf_body.word_wrap = True
        for i, line in enumerate(body):
            if i == 0:
                p = tf_body.paragraphs[0]
            else:
                p = tf_body.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.space_after = Pt(4)
            if line.strip().startswith("•") or (line.strip() and not line.startswith(" ")):
                p.level = 0
            elif line.strip().startswith("  "):
                p.level = 1

    # Speaker note
    if slide_spec.get("note"):
        slide.notes_slide.notes_text_frame.text = slide_spec["note"]

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        add_slide(prs, spec)

    out_path = "Alabama_Tree_Cover_GEE_Presentation.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
